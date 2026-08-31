import json
import sys
from types import ModuleType

MAX_SLIDES = 500
MAX_SHAPES_PER_SLIDE = 10_000
MAX_TABLE_CELLS_PER_SLIDE = 50_000
MAX_SLIDE_CODE_POINTS = 100_000
MAX_DOCUMENT_CODE_POINTS = 2_000_000


class BlockedPillow:
    @staticmethod
    def open(_stream):
        raise ValueError("image_decoding_unsupported")

    @staticmethod
    def truetype(_font_path, _point_size):
        raise ValueError("font_rendering_unsupported")


def install_blocked_pillow():
    module = ModuleType("PIL")
    module.Image = BlockedPillow
    module.ImageFont = BlockedPillow
    sys.modules["PIL"] = module


install_blocked_pillow()
from pptx import Presentation  # noqa: E402


def shape_text(shape, state):
    if getattr(shape, "has_text_frame", False):
        return shape.text.strip()
    if getattr(shape, "has_table", False):
        chunks: list[str] = []
        for row in shape.table.rows:
            for cell in row.cells:
                state["cells"] += 1
                if state["cells"] > MAX_TABLE_CELLS_PER_SLIDE:
                    raise ValueError("table_cell_limit")
                content = cell.text.strip()
                if content:
                    chunks.append(content)
        return "\n".join(chunks)
    return ""


def slide_text(slide):
    chunks: list[str] = []
    state = {"shapes": 0, "cells": 0}
    queue = list(slide.shapes)
    cursor = 0
    while cursor < len(queue):
        shape = queue[cursor]
        cursor += 1
        state["shapes"] += 1
        if state["shapes"] > MAX_SHAPES_PER_SLIDE:
            raise ValueError("shape_limit")
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            queue.extend(list(nested))
            if len(queue) + state["shapes"] > MAX_SHAPES_PER_SLIDE:
                raise ValueError("shape_limit")
        content = shape_text(shape, state)
        if content:
            chunks.append(content)
    text = "\n".join(chunks)
    if len(text) > MAX_SLIDE_CODE_POINTS:
        raise ValueError("slide_limit")
    return text


def parse_pptx(path):
    presentation = Presentation(path)
    sources: list[dict[str, object]] = []
    total_points = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index > MAX_SLIDES:
            raise ValueError("slide_count_limit")
        content = slide_text(slide)
        total_points += len(content)
        if total_points > MAX_DOCUMENT_CODE_POINTS:
            raise ValueError("document_limit")
        if content:
            sources.append({"slide": slide_index, "content": content})
    if not sources:
        raise ValueError("no_text")
    return json.dumps(
        {"schema_version": "1", "format": "pptx", "sources": sources},
        ensure_ascii=False,
        separators=(",", ":"),
    )


parse_pptx("/tmp/aethelgard-source.pptx")
